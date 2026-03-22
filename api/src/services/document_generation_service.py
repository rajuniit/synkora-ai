"""Document Generation Service for creating PDF, Google Docs, Sheets, and PowerPoint files."""

import io
import logging
import re
from datetime import UTC, datetime
from typing import Any
from uuid import UUID

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    HRFlowable,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from sqlalchemy.ext.asyncio import AsyncSession

from src.services.storage.s3_storage import S3StorageService

logger = logging.getLogger(__name__)


class DocumentGenerationService:
    """Service for generating documents in various formats."""

    def __init__(self, db: AsyncSession, tenant_id: UUID):
        """Initialize document generation service.

        Args:
            db: Database session
            tenant_id: Tenant ID
        """
        self.db = db
        self.tenant_id = tenant_id
        self.storage = S3StorageService()

    def _inline_markup(self, text: str) -> str:
        """Convert inline markdown to ReportLab XML markup with proper escaping."""
        # Markers that won't appear in normal text
        BOLD_O, BOLD_C = "\x00b\x00", "\x01b\x01"
        ITALIC_O, ITALIC_C = "\x00i\x00", "\x01i\x01"
        CODE_O, CODE_C = "\x00c\x00", "\x01c\x01"

        # Strip HTML tags from LLM output
        text = re.sub(r"<strong>(.*?)</strong>", r"**\1**", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"<em>(.*?)</em>", r"*\1*", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"<[^>]+>", "", text)

        # Apply markdown → markers (order matters: ** before *)
        text = re.sub(r"\*\*(.+?)\*\*", lambda m: BOLD_O + m.group(1) + BOLD_C, text)
        text = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", lambda m: ITALIC_O + m.group(1) + ITALIC_C, text)
        text = re.sub(r"`(.+?)`", lambda m: CODE_O + m.group(1) + CODE_C, text)

        # Escape XML special chars in plain text segments (between markers)
        all_markers = {BOLD_O, BOLD_C, ITALIC_O, ITALIC_C, CODE_O, CODE_C}
        parts = re.split(r"(\x00[bic]\x00|\x01[bic]\x01)", text)
        escaped = []
        for part in parts:
            if part in all_markers:
                escaped.append(part)
            else:
                escaped.append(part.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))
        text = "".join(escaped)

        # Replace markers with ReportLab XML tags
        text = (
            text.replace(BOLD_O, "<b>")
            .replace(BOLD_C, "</b>")
            .replace(ITALIC_O, "<i>")
            .replace(ITALIC_C, "</i>")
            .replace(CODE_O, '<font name="Courier" fontSize="9">')
            .replace(CODE_C, "</font>")
        )
        return text

    def _build_table_element(self, table_lines: list[str], table_cell_style: Any, table_header_style: Any) -> Any:
        """Build a ReportLab Table from markdown table lines."""
        rows = []
        for line in table_lines:
            cells = [c.strip() for c in line.strip("|").split("|")]
            rows.append(cells)

        if not rows:
            return None

        col_count = max(len(row) for row in rows)

        # Build styled rows
        styled_rows = []
        for r_idx, row in enumerate(rows):
            style = table_header_style if r_idx == 0 else table_cell_style
            styled_row = []
            for c_idx in range(col_count):
                cell_text = row[c_idx] if c_idx < len(row) else ""
                styled_row.append(Paragraph(self._inline_markup(cell_text), style))
            styled_rows.append(styled_row)

        # Calculate column widths (equal distribution within page width)
        available_width = letter[0] - 144  # page width minus margins
        col_width = available_width / col_count

        table = Table(styled_rows, colWidths=[col_width] * col_count, repeatRows=1)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f3f4f6")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d1d5db")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f9fafb")]),
                ]
            )
        )
        return table

    def _parse_markdown_to_elements(self, content: str, styles_dict: dict) -> list:
        """Parse markdown content into a list of ReportLab flowable elements."""
        elements = []
        lines = content.split("\n")
        i = 0

        while i < len(lines):
            line = lines[i]
            stripped = line.strip()

            # Code block
            if stripped.startswith("```"):
                i += 1
                code_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    code_lines.append(lines[i])
                    i += 1
                if code_lines:
                    elements.append(Preformatted("\n".join(code_lines), styles_dict["code"]))
                    elements.append(Spacer(1, 6))

            # Heading 1 (not ##)
            elif re.match(r"^# [^#]", stripped):
                text = stripped[2:].strip()
                elements.append(Paragraph(self._inline_markup(text), styles_dict["h1"]))
                elements.append(Spacer(1, 8))

            # Heading 2 (not ###)
            elif re.match(r"^## [^#]", stripped):
                text = stripped[3:].strip()
                elements.append(Paragraph(self._inline_markup(text), styles_dict["h2"]))
                elements.append(Spacer(1, 6))

            # Heading 3+
            elif re.match(r"^#{3,6} ", stripped):
                text = re.sub(r"^#{3,6} ", "", stripped)
                elements.append(Paragraph(self._inline_markup(text), styles_dict["h3"]))
                elements.append(Spacer(1, 4))

            # Horizontal rule
            elif re.match(r"^-{3,}$|^\*{3,}$|^_{3,}$", stripped):
                elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#d1d5db")))
                elements.append(Spacer(1, 6))

            # Markdown table row
            elif stripped.startswith("|"):
                table_lines = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    row = lines[i].strip()
                    # Skip separator rows like |---|---|
                    if not re.match(r"^\|[\s\-\|:]+\|$", row):
                        table_lines.append(row)
                    i += 1
                i -= 1  # will be incremented at end of loop
                if table_lines:
                    tbl = self._build_table_element(table_lines, styles_dict["table_cell"], styles_dict["table_header"])
                    if tbl:
                        elements.append(tbl)
                        elements.append(Spacer(1, 10))

            # Bullet list (- item or * item)
            elif re.match(r"^[-*]\s+", stripped):
                while i < len(lines) and re.match(r"^[-*]\s+", lines[i].strip()):
                    item = re.sub(r"^[-*]\s+", "", lines[i].strip())
                    elements.append(Paragraph("• " + self._inline_markup(item), styles_dict["bullet"]))
                    i += 1
                i -= 1
                elements.append(Spacer(1, 4))

            # Numbered list
            elif re.match(r"^\d+\.\s+", stripped):
                while i < len(lines) and re.match(r"^\d+\.\s+", lines[i].strip()):
                    item = re.sub(r"^\d+\.\s+", "", lines[i].strip())
                    num = re.match(r"^(\d+)\.", lines[i].strip()).group(1)
                    elements.append(Paragraph(f"{num}. " + self._inline_markup(item), styles_dict["bullet"]))
                    i += 1
                i -= 1
                elements.append(Spacer(1, 4))

            # Blockquote
            elif stripped.startswith("> "):
                text = stripped[2:]
                elements.append(Paragraph(f"<i>{self._inline_markup(text)}</i>", styles_dict["quote"]))
                elements.append(Spacer(1, 4))

            # Empty line → small spacer
            elif stripped == "":
                if elements and not isinstance(elements[-1], Spacer):
                    elements.append(Spacer(1, 6))

            # Regular paragraph text
            else:
                # Accumulate consecutive regular lines into one paragraph
                para_lines = [stripped]
                while i + 1 < len(lines):
                    next_s = lines[i + 1].strip()
                    if (
                        not next_s
                        or next_s.startswith("#")
                        or next_s.startswith("|")
                        or next_s.startswith("```")
                        or re.match(r"^[-*]\s+", next_s)
                        or re.match(r"^\d+\.\s+", next_s)
                        or next_s.startswith("> ")
                        or re.match(r"^-{3,}$|^\*{3,}$", next_s)
                    ):
                        break
                    i += 1
                    para_lines.append(lines[i].strip())

                text = " ".join(p for p in para_lines if p)
                if text:
                    elements.append(Paragraph(self._inline_markup(text), styles_dict["body"]))
                    elements.append(Spacer(1, 6))

            i += 1

        return elements

    async def generate_pdf(
        self,
        content: str,
        filename: str | None = None,
        title: str = "Document",
        author: str = "Synkora",
        include_metadata: bool = True,
    ) -> dict[str, Any]:
        """Generate PDF document from markdown or plain text content."""
        import os

        try:
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(
                pdf_buffer,
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=36,
            )

            styles = getSampleStyleSheet()
            styles_dict = {
                "h1": ParagraphStyle(
                    "H1",
                    parent=styles["Heading1"],
                    fontSize=20,
                    textColor=colors.HexColor("#111827"),
                    spaceBefore=14,
                    spaceAfter=6,
                    fontName="Helvetica-Bold",
                ),
                "h2": ParagraphStyle(
                    "H2",
                    parent=styles["Heading2"],
                    fontSize=16,
                    textColor=colors.HexColor("#1f2937"),
                    spaceBefore=12,
                    spaceAfter=4,
                    fontName="Helvetica-Bold",
                ),
                "h3": ParagraphStyle(
                    "H3",
                    parent=styles["Heading3"],
                    fontSize=13,
                    textColor=colors.HexColor("#374151"),
                    spaceBefore=8,
                    spaceAfter=3,
                    fontName="Helvetica-Bold",
                ),
                "body": ParagraphStyle(
                    "Body",
                    parent=styles["BodyText"],
                    fontSize=11,
                    textColor=colors.HexColor("#374151"),
                    spaceAfter=6,
                    leading=16,
                ),
                "bullet": ParagraphStyle(
                    "Bullet",
                    parent=styles["BodyText"],
                    fontSize=11,
                    textColor=colors.HexColor("#374151"),
                    leftIndent=20,
                    spaceAfter=3,
                    leading=15,
                ),
                "quote": ParagraphStyle(
                    "Quote",
                    parent=styles["BodyText"],
                    fontSize=11,
                    textColor=colors.HexColor("#6b7280"),
                    leftIndent=24,
                    spaceAfter=4,
                    leading=15,
                    borderPadding=(2, 2, 2, 8),
                ),
                "code": ParagraphStyle(
                    "Code",
                    parent=styles["Code"],
                    fontSize=9,
                    fontName="Courier",
                    backColor=colors.HexColor("#f3f4f6"),
                    leftIndent=12,
                    rightIndent=12,
                    spaceAfter=4,
                ),
                "table_header": ParagraphStyle(
                    "TH",
                    parent=styles["BodyText"],
                    fontSize=10,
                    fontName="Helvetica-Bold",
                    textColor=colors.HexColor("#111827"),
                ),
                "table_cell": ParagraphStyle(
                    "TD",
                    parent=styles["BodyText"],
                    fontSize=10,
                    textColor=colors.HexColor("#374151"),
                    leading=13,
                ),
            }

            title_style = ParagraphStyle(
                "DocTitle",
                parent=styles["Heading1"],
                fontSize=24,
                textColor=colors.HexColor("#1a1a1a"),
                spaceAfter=8,
                alignment=1,
                fontName="Helvetica-Bold",
            )
            meta_style = ParagraphStyle(
                "Meta",
                parent=styles["BodyText"],
                fontSize=9,
                textColor=colors.HexColor("#9ca3af"),
                alignment=1,
            )

            elements: list = []
            elements.append(Paragraph(title, title_style))
            if include_metadata:
                elements.append(
                    Paragraph(
                        f"Generated: {datetime.now(UTC).strftime('%Y-%m-%d %H:%M UTC')} &nbsp;|&nbsp; Author: {author}",
                        meta_style,
                    )
                )
            elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#e5e7eb")))
            elements.append(Spacer(1, 12))

            if content:
                elements.extend(self._parse_markdown_to_elements(content, styles_dict))
            else:
                elements.append(Paragraph("<i>No content provided</i>", styles_dict["body"]))

            doc.build(elements)
            pdf_content = pdf_buffer.getvalue()

            if not filename:
                timestamp = datetime.now(UTC).strftime("%Y%m%d_%H%M%S")
                filename = f"{title.replace(' ', '_')}_{timestamp}.pdf"
            elif not filename.endswith(".pdf"):
                filename = f"{filename}.pdf"

            now = datetime.now(UTC)
            file_path = f"documents/{self.tenant_id}/{now.strftime('%Y/%m/%d')}/{filename}"

            logger.info(f"Uploading PDF to S3: {file_path}")
            result = self.storage.upload_file(file_content=pdf_content, key=file_path, content_type="application/pdf")
            logger.info(f"PDF uploaded: {result['key']}")

            download_url = self.storage.generate_presigned_url(result["key"], expiration=86400 * 7)

            # Replace internal MinIO endpoint with public endpoint in presigned URL
            internal_endpoint = os.getenv("AWS_ENDPOINT_URL", "")
            public_endpoint = os.getenv("AWS_PUBLIC_ENDPOINT_URL", "")
            if internal_endpoint and public_endpoint and internal_endpoint != public_endpoint:
                download_url = download_url.replace(internal_endpoint, public_endpoint)

            logger.info(f"PDF ready — download: {download_url[:80]}...")

            return {
                "success": True,
                "format": "pdf",
                "file_path": result["key"],
                "file_name": filename,
                "file_size": len(pdf_content),
                "title": title,
                "pages": doc.page,
                "download_url": download_url,
                "message": f"📄 PDF generated successfully!\n\nDownload: {download_url}",
            }

        except Exception as e:
            logger.error(f"Failed to generate PDF: {e}", exc_info=True)
            return {"success": False, "message": f"PDF generation failed: {str(e)}", "error": str(e)}

    async def generate_powerpoint(
        self, slides_content: list[dict[str, Any]], filename: str | None = None, title: str = "Presentation"
    ) -> dict[str, Any]:
        """Generate PowerPoint presentation.

        Args:
            slides_content: List of slide dictionaries with 'title' and 'content' or 'bullet_points'
            filename: Custom filename (default: auto-generated)
            title: Presentation title

        Returns:
            Dict with generation info and file path
        """
        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = f"Generated on {datetime.now(UTC).strftime('%B %d, %Y')}"

            # Add content slides
            for slide_data in slides_content:
                slide_title = slide_data.get("title", "Slide")
                content = slide_data.get("content", "")
                bullet_points = slide_data.get("bullet_points", [])

                # Use bullet layout if we have bullet points
                if bullet_points:
                    bullet_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    slide.shapes.title.text = slide_title

                    text_frame = slide.shapes.placeholders[1].text_frame
                    text_frame.clear()

                    for point in bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                else:
                    # Use blank layout for custom content
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)

                    # Add title
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                    title_frame = title_box.text_frame
                    title_para = title_frame.add_paragraph()
                    title_para.text = slide_title
                    title_para.font.size = Pt(32)
                    title_para.font.bold = True

                    # Add content
                    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
                    content_frame = content_box.text_frame
                    content_para = content_frame.add_paragraph()
                    content_para.text = content
                    content_para.font.size = Pt(18)

            # Save to buffer
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_content = pptx_buffer.getvalue()

            # Generate filename
            if not filename:
                timestamp = datetime.now(UTC).strftime("%Y%m%d_%H%M%S")
                filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
            elif not filename.endswith(".pptx"):
                filename = f"{filename}.pptx"

            # Upload to S3
            now = datetime.now(UTC)
            date_path = now.strftime("%Y/%m/%d")
            file_path = f"documents/{self.tenant_id}/{date_path}/{filename}"

            logger.info(f"Uploading PowerPoint to S3: {file_path}")
            result = self.storage.upload_file(
                file_content=pptx_content,
                key=file_path,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

            logger.info(f"PowerPoint uploaded successfully to: {result['key']}")

            import os

            download_url = self.storage.generate_presigned_url(result["key"], expiration=86400 * 7)

            # Replace internal MinIO endpoint with public endpoint in presigned URL
            internal_endpoint = os.getenv("AWS_ENDPOINT_URL", "")
            public_endpoint = os.getenv("AWS_PUBLIC_ENDPOINT_URL", "")
            if internal_endpoint and public_endpoint and internal_endpoint != public_endpoint:
                download_url = download_url.replace(internal_endpoint, public_endpoint)

            logger.info(f"PowerPoint ready — download: {download_url[:80]}...")

            return {
                "success": True,
                "format": "powerpoint",
                "file_path": result["key"],
                "file_name": filename,
                "file_size": len(pptx_content),
                "title": title,
                "slides": len(prs.slides),
                "download_url": download_url,
                "message": f"📊 PowerPoint generated successfully!\n\nDownload: {download_url}",
            }

        except Exception as e:
            logger.error(f"Failed to generate PowerPoint: {e}", exc_info=True)
            return {"success": False, "message": f"PowerPoint generation failed: {str(e)}", "error": str(e)}

    async def generate_google_doc(
        self, content: str, title: str = "Document", share_with_emails: list[str] | None = None
    ) -> dict[str, Any]:
        """Generate Google Doc (requires Google API credentials).

        Args:
            content: Document content
            title: Document title
            share_with_emails: List of emails to share document with

        Returns:
            Dict with generation info and document URL
        """
        try:
            # Note: This requires OAuth tokens to be stored in the database
            # For now, return a message about setup requirement
            return {
                "success": False,
                "message": "Google Docs integration requires OAuth setup. Please configure Google API credentials first.",
                "setup_required": True,
                "instructions": "1. Enable Google Docs API in Google Cloud Console\n2. Configure OAuth credentials\n3. Grant access to create documents",
            }

        except Exception as e:
            logger.error(f"Failed to generate Google Doc: {e}", exc_info=True)
            return {"success": False, "message": f"Google Doc generation failed: {str(e)}", "error": str(e)}

    async def generate_google_sheet(
        self,
        data: list[list[Any]],
        title: str = "Spreadsheet",
        sheet_name: str = "Sheet1",
        share_with_emails: list[str] | None = None,
    ) -> dict[str, Any]:
        """Generate Google Sheet (requires Google API credentials).

        Args:
            data: 2D array of data for the sheet
            title: Spreadsheet title
            sheet_name: Sheet name
            share_with_emails: List of emails to share sheet with

        Returns:
            Dict with generation info and sheet URL
        """
        try:
            # Note: This requires OAuth tokens to be stored in the database
            # For now, return a message about setup requirement
            return {
                "success": False,
                "message": "Google Sheets integration requires OAuth setup. Please configure Google API credentials first.",
                "setup_required": True,
                "instructions": "1. Enable Google Sheets API in Google Cloud Console\n2. Configure OAuth credentials\n3. Grant access to create spreadsheets",
            }

        except Exception as e:
            logger.error(f"Failed to generate Google Sheet: {e}", exc_info=True)
            return {"success": False, "message": f"Google Sheet generation failed: {str(e)}", "error": str(e)}
